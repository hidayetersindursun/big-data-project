# -*- coding: utf-8 -*-
"""GıdaRadar — Final sunum üreticisi (.pptx + konusma_metni.md).

Slaytlar İngilizce, konuşma metni (notlar) Türkçe.
Kullanım:  python presentation/build_presentation.py
"""
import os
import re
import sys

sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")

# ---- renk paleti -----------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
ORANGE = RGBColor(0xE7, 0x6F, 0x51)
GOLD = RGBColor(0xE9, 0xC4, 0x6A)
GREY = RGBColor(0x5B, 0x66, 0x77)
LIGHT = RGBColor(0xF1, 0xF3, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2A, 0x35)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

NOTES = []          # konusma_metni.md icin (baslik, metin)
_slide_no = [0]


# ---- yardimcilar -----------------------------------------------------------
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _solid(r, bg)
    r.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size, color=DARK, bold=False, first=False,
         align=PP_ALIGN.LEFT, italic=False, space_after=6, bullet=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if bullet is not None:
        run = p.add_run()
        run.text = bullet + "  "
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = TEAL
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def title_bar(slide, title, kicker=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.18))
    _solid(bar, NAVY)
    bar.shadow.inherit = False
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.18),
                                    SW, Inches(0.07))
    _solid(accent, TEAL)
    accent.shadow.inherit = False
    tf = textbox(slide, 0.55, 0.06, 11.0, 1.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if kicker:
        para(tf, kicker.upper(), 13, GOLD, bold=True, first=True, space_after=2)
        para(tf, title, 29, WHITE, bold=True)
    else:
        para(tf, title, 30, WHITE, bold=True, first=True)


def footer(slide):
    _slide_no[0] += 1
    tf = textbox(slide, 0.55, 7.04, 9.0, 0.4)
    para(tf, "GıdaRadar  ·  Türkiye Food Supply Chain Transparency Engine",
         10, GREY, first=True)
    tn = textbox(slide, 12.2, 7.04, 0.9, 0.4)
    para(tn, str(_slide_no[0]), 10, GREY, first=True, align=PP_ALIGN.RIGHT)


def bullets(slide, items, x, y, w, h, size=20, gap=10):
    """items: (text, level) ya da text."""
    tf = textbox(slide, x, y, w, h)
    first = True
    for it in items:
        text, lvl = (it if isinstance(it, tuple) else (it, 0))
        if lvl == 0:
            para(tf, text, size, DARK, bold=False, first=first,
                 bullet="▸", space_after=gap)
        elif lvl == 1:
            p = para(tf, text, size - 4, GREY, first=first,
                     bullet="–", space_after=gap - 3)
            p.level = 1
        else:  # baslik satiri
            para(tf, text, size, NAVY, bold=True, first=first, space_after=gap)
        first = False
    return tf


def chip(slide, x, y, w, h, text, fill, fg=WHITE, size=14, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _solid(box, fill)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, text, size, fg, bold=True, first=True, align=PP_ALIGN.CENTER)
    return box


def arrow(slide, x, y, w, h=0.32, color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    _solid(a, color)
    a.shadow.inherit = False


def picture(slide, path, x, y, w=None, h=None):
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(os.path.join(ASSETS, path),
                                    Inches(x), Inches(y), **kw)


def caption(slide, text, x, y, w):
    tf = textbox(slide, x, y, w, 0.5)
    para(tf, text, 12, GREY, italic=True, first=True)


def notes(slide, title, text):
    # Slayt numarasi otomatik turetilir; gelen baslikta numara varsa atilir.
    label = re.sub(r"^\d+\s*[—–-]\s*", "", title)
    full = f"{len(NOTES) + 1} — {label}"
    slide.notes_slide.notes_text_frame.text = text
    NOTES.append((full, text))


def result_slide(kicker, title, img, findings, takeaway):
    s = new_slide()
    title_bar(s, title, kicker)
    pic = picture(s, img, 0.5, 1.55, w=7.7)
    # bulgular paneli
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(8.45), Inches(1.55),
                               Inches(4.35), Inches(4.05))
    _solid(panel, LIGHT)
    panel.shadow.inherit = False
    tf = textbox(s, 8.7, 1.75, 3.9, 3.7)
    para(tf, "KEY FINDINGS", 14, TEAL, bold=True, first=True, space_after=10)
    for f in findings:
        para(tf, f, 16, DARK, bullet="▸", space_after=9)
    # takeaway seridi
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5),
                              Inches(5.85), Inches(12.3), Inches(0.92))
    _solid(band, NAVY)
    band.shadow.inherit = False
    tf2 = band.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf2, "INSIGHT   " + takeaway, 16, WHITE, bold=True, first=True,
         align=PP_ALIGN.CENTER)
    footer(s)
    return s


# ===========================================================================
# SLAYT 1 — TITLE
# ===========================================================================
s = new_slide(NAVY)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.85), SW, Inches(0.09))
_solid(band, TEAL)
band.shadow.inherit = False
tf = textbox(s, 1.0, 0.95, 11.3, 1.0)
para(tf, "BIG DATA — FINAL PROJECT", 18, GOLD, bold=True, first=True,
     align=PP_ALIGN.CENTER)
tf = textbox(s, 0.7, 1.5, 11.9, 1.4)
para(tf, "GıdaRadar", 60, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = textbox(s, 0.7, 3.05, 11.9, 1.5)
para(tf, "Türkiye Food Supply Chain Transparency", 30, WHITE, bold=True,
     first=True, align=PP_ALIGN.CENTER)
para(tf, "& Spatial Margin Analysis Engine", 30, TEAL, bold=True,
     align=PP_ALIGN.CENTER)
tf = textbox(s, 0.7, 4.7, 11.9, 0.9)
para(tf, "Wholesale (Hal) → Retail margins · Rockets & Feathers ·"
         " Shock propagation · Spatial speculation", 17, RGBColor(0xC9, 0xD2,
         0xE0), first=True, align=PP_ALIGN.CENTER)
tf = textbox(s, 0.7, 5.7, 11.9, 0.6)
para(tf, "Azmi Yağlı   ·   Abdullah Zengin   ·   Hidayet Ersin Dursun",
     20, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = textbox(s, 0.7, 6.55, 11.9, 0.5)
para(tf, "Bu çalışmada yapay zeka araçlarından (AI assist) yararlanılmıştır.",
     12, RGBColor(0x9A, 0xA6, 0xB8), italic=True, first=True,
     align=PP_ALIGN.CENTER)
notes(s, "1 — Başlık",
      "Herkese merhaba. Biz Azmi Yağlı, Abdullah Zengin ve Hidayet Ersin "
      "Dursun olarak Big Data dersi final projemiz GıdaRadar'ı sunuyoruz. "
      "GıdaRadar, Türkiye'de gıda tedarik zincirinin şeffaflığını ölçen bir "
      "büyük veri motoru. Temel sorumuz şu: bir ürün haldeki toptan "
      "fiyatından market rafına gelene kadar fiyatı nasıl değişiyor, bu marj "
      "nerede, ne zaman, neden açılıyor? Projede toptan ile perakende "
      "arasındaki marjı, asimetrik fiyat geçişini yani Rockets and Feathers "
      "etkisini, hava şoklarının fiyata yansıma hızını ve şehirler arası "
      "fiyat eşitsizliğini analiz ettik. Sunum boyunca önce problemi ve "
      "amacı, sonra veri toplama ve mimariyi, ardından AWS üzerinde "
      "çalıştırdığımız pipeline'ı ve en sonunda gerçek analiz sonuçlarımızı "
      "anlatacağız. Belirtmek isteriz: çalışmamızda yapay zeka araçlarından "
      "yararlandık, bunu kaynakça slaytında da not ettik.")
footer(s)

# ===========================================================================
# SLAYT 2 — TEAM
# ===========================================================================
s = new_slide()
title_bar(s, "The Team", "Who we are")
members = [
    ("Azmi Yağlı", "Data Engineering",
     "Ingestion & scraping, NiFi/Kafka, Bronze layer, S3 partitioning"),
    ("Abdullah Zengin", "Platform & Pipeline",
     "Medallion architecture, Spark Silver layer, EMR orchestration"),
    ("Hidayet Ersin Dursun", "Data Science",
     "Entity resolution, Gold analytics, forecasting & dashboards"),
]
x = 0.7
for name, role, desc in members:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                              Inches(1.8), Inches(3.85), Inches(3.5))
    _solid(card, LIGHT)
    card.shadow.inherit = False
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                              Inches(1.8), Inches(3.85), Inches(0.95))
    _solid(head, NAVY)
    head.shadow.inherit = False
    tfh = head.text_frame
    tfh.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tfh, name, 18, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
    tf = textbox(s, x + 0.25, 3.0, 3.35, 2.2)
    para(tf, role, 17, TEAL, bold=True, first=True, space_after=10,
         align=PP_ALIGN.CENTER)
    para(tf, desc, 15, DARK, align=PP_ALIGN.CENTER)
    x += 4.07
tf = textbox(s, 0.7, 5.7, 11.9, 0.9)
para(tf, "A 3-person team carrying both the Data Engineering and the "
         "Data Science hats end-to-end.", 16, GREY, italic=True, first=True,
     align=PP_ALIGN.CENTER)
notes(s, "2 — Ekip",
      "Ekibimiz üç kişiden oluşuyor ve projenin hem veri mühendisliği hem de "
      "veri bilimi tarafını birlikte taşıdık. Azmi veri toplama ve scraping "
      "tarafına, NiFi ve Kafka ile ham verinin S3 Bronze katmanına "
      "akmasına odaklandı. Abdullah medallion mimarisini, Spark ile Silver "
      "katmanını ve EMR üzerindeki pipeline orkestrasyonunu üstlendi. Ben "
      "Hidayet ise entity resolution, Gold analiz tabloları, tahminleme ve "
      "dashboard tarafını yürüttüm. Pratikte sınırlar keskin değildi; üç "
      "kişi olduğumuz için herkes her katmana dokundu. Bu yüzden sunumu da "
      "uçtan uca, ortak bir anlatı olarak hazırladık.")
footer(s)

# ===========================================================================
# SLAYT 3 — PROBLEM
# ===========================================================================
s = new_slide()
title_bar(s, "Problem Statement", "Why this matters")
intro = textbox(s, 0.7, 1.5, 12.0, 0.8)
para(intro, "A tomato leaves the wholesale market (hal) at ~15 ₺/kg and "
            "reaches the shelf at 30–40 ₺/kg. Where does the gap come from "
            "— and is it fair?", 18, DARK, bold=True, first=True)
cards = [
    ("The Black Box", ORANGE,
     "No public dataset links wholesale prices to retail shelf prices. "
     "The margin is structurally invisible to regulators and consumers."),
    ("Rockets & Feathers", NAVY,
     "When wholesale costs rise, shelf prices jump like a rocket. When "
     "costs fall, they drift down like a feather — asymmetric transmission."),
    ("Spatial Speculation", TEAL,
     "Identical wholesale entry price, very different shelf prices across "
     "provinces — location-based pricing power goes unmeasured."),
]
x = 0.7
for t, c, d in cards:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                              Inches(2.55), Inches(3.85), Inches(3.05))
    _solid(card, LIGHT)
    card.shadow.inherit = False
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                               Inches(2.55), Inches(3.85), Inches(0.18))
    _solid(strip, c)
    strip.shadow.inherit = False
    tf = textbox(s, x + 0.28, 2.85, 3.3, 2.6)
    para(tf, t, 19, c, bold=True, first=True, space_after=10)
    para(tf, d, 15, DARK)
    x += 4.07
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7),
                          Inches(5.85), Inches(11.93), Inches(0.85))
_solid(band, NAVY)
band.shadow.inherit = False
tfb = band.text_frame
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfb, "Food inflation is a national priority — yet the margin layer "
          "between farm and shelf has never been measured at scale.", 15,
     WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "3 — Problem",
      "Projemizin çıkış noktası çok somut bir gözlem. Bir domates haldeki "
      "toptan satıştan yaklaşık 15 lira civarında çıkıyor ama market rafına "
      "30, 40 liraya ulaşıyor. Bu farkın nereden geldiğini kimse net "
      "ölçemiyor. Üç temel problem var. Birincisi kara kutu: toptan fiyatı "
      "ile perakende rafiyatını birbirine bağlayan kamuya açık tek bir veri "
      "seti yok, dolayısıyla marj yapısal olarak görünmez. İkincisi Rockets "
      "and Feathers etkisi: maliyet artınca raf fiyatı roket gibi fırlıyor, "
      "maliyet düşünce tüy gibi yavaşça iniyor; yani asimetrik bir geçiş "
      "var. Üçüncüsü mekânsal spekülasyon: aynı hal giriş fiyatına sahip "
      "ürün farklı illerde çok farklı fiyatlanıyor ve bu lokasyon bazlı "
      "fiyatlama gücü hiç ölçülmüyor. Gıda enflasyonu ülke gündeminin en "
      "tepesinde, ama tarla ile raf arasındaki bu marj katmanı büyük "
      "ölçekte hiç incelenmemiş. İşte bu boşluğu doldurmak istedik.")
footer(s)

# ===========================================================================
# SLAYT 4 — REAL-WORLD CASE (gerçek ceza haberi)
# ===========================================================================
s = new_slide()
title_bar(s, "A Real Case — Not Hypothetical", "In the news · March 2026")
picture(s, "file_market_ceza.png", 0.55, 1.6, w=7.65)
caption(s, "Web search result — \"file market ceza\", March 2026.",
        0.6, 5.02, 7.6)
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45),
                           Inches(1.6), Inches(4.35), Inches(3.42))
_solid(panel, LIGHT)
panel.shadow.inherit = False
tf = textbox(s, 8.72, 1.8, 3.85, 3.15)
para(tf, "WHAT THE HEADLINE SAYS", 14, TEAL, bold=True, first=True,
     space_after=11)
for t in ["A national market chain was fined ₺1.8 M for excessive pricing.",
          "A pepper bought at ₺50 from the hal was put on the shelf "
          "at ₺190.",
          "That is a ~280% markup on a single staple vegetable.",
          "The regulator caught it once — through a manual inspection."]:
    para(tf, t, 14.5, DARK, bullet="▸", space_after=10)
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55),
                          Inches(5.5), Inches(12.25), Inches(1.2))
_solid(band, NAVY)
band.shadow.inherit = False
tfb = band.text_frame
tfb.word_wrap = True
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfb, "GıdaRadar measures exactly this margin — automatically, every "
          "day, across 81 provinces and every chain. One manual catch "
          "becomes systematic, scalable detection.", 16, WHITE, bold=True,
     first=True, align=PP_ALIGN.CENTER)
notes(s, "Real Case — Gerçek Vaka",
      "Problemi soyut bırakmayalım — bu gerçek bir haber. Mart 2026'da "
      "Ticaret Bakanlığı, ulusal bir market zinciri olan File Market'e "
      "fahiş fiyat nedeniyle 1 milyon 806 bin lira idari para cezası "
      "kesti. Denetimde tespit edilen şey tam da bizim projemizin ölçtüğü "
      "şey: haldeki toptan fiyatı 50 lira olan biber, rafa 190 liraya "
      "konmuş — yani tek bir temel sebzede yüzde 280 kâr marjı. Burada iki "
      "nokta kritik. Birincisi: bu vaka, GıdaRadar'ın ürettiği marj "
      "sayılarının gerçek dünyada birebir karşılığı olduğunu kanıtlıyor; "
      "marj sonuçları slaytında göreceğimiz uzun sağ kuyruk, işte tam "
      "olarak bu tür uç vakalardan oluşuyor. İkincisi: düzenleyici bu "
      "vakayı tek tek, manuel bir denetimle yakaladı. GıdaRadar ise aynı "
      "örüntüyü otomatik olarak, her gün, 81 ilde ve tüm market "
      "zincirlerinde tespit ediyor. Yani tek bir manuel yakalamayı "
      "sistematik ve ölçeklenebilir bir erken uyarıya dönüştürüyoruz. Bu "
      "slayt hem problemin gerçekliğini hem de projenin somut iş değerini "
      "tek karede gösteriyor.")
footer(s)

# ===========================================================================
# SLAYT 5 — GOAL
# ===========================================================================
s = new_slide()
title_bar(s, "Project Goal", "From data to decisions")
vis = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5),
                         Inches(11.93), Inches(1.15))
_solid(vis, TEAL)
vis.shadow.inherit = False
tfv = vis.text_frame
tfv.vertical_anchor = MSO_ANCHOR.MIDDLE
tfv.word_wrap = True
para(tfv, "Build a Decision-Support System & Early-Warning Radar that makes "
          "the food supply chain margin measurable, comparable and "
          "predictable.", 18, WHITE, bold=True, first=True,
     align=PP_ALIGN.CENTER)
outs = [
    ("Margin X-Ray", "Daily hal→retail margin by product, chain & province"),
    ("Asymmetry Score", "Quantify Rockets & Feathers pricing per product"),
    ("Shock Radar", "Measure how fast weather shocks hit shelf prices"),
    ("Forecasting", "Prophet-based price outlook & structural-break detection"),
]
x = 0.7
for t, d in outs:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                              Inches(2.95), Inches(2.86), Inches(2.5))
    _solid(card, LIGHT)
    card.shadow.inherit = False
    tf = textbox(s, x + 0.2, 3.2, 2.5, 2.1)
    para(tf, t, 17, NAVY, bold=True, first=True, space_after=8,
         align=PP_ALIGN.CENTER)
    para(tf, d, 14, DARK, align=PP_ALIGN.CENTER)
    x += 3.04
tf = textbox(s, 0.7, 5.75, 11.93, 1.0)
para(tf, "Primary stakeholders:  Competition Authority · Ministry of Trade · "
         "Central Bank (food-inflation analysis) · consumers",
     16, GREY, bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "4 — Amaç",
      "Amacımız tek cümleyle şu: gıda tedarik zincirindeki marjı ölçülebilir, "
      "karşılaştırılabilir ve tahmin edilebilir hale getiren bir karar "
      "destek sistemi ve erken uyarı radarı kurmak. Bunu dört somut çıktıyla "
      "tarif ediyoruz. Marj röntgeni: ürün, market zinciri ve il bazında "
      "günlük hal-perakende marjı. Asimetri skoru: her ürün için Rockets and "
      "Feathers davranışını sayısallaştırmak. Şok radarı: hava şoklarının raf "
      "fiyatına kaç günde yansıdığını ölçmek. Ve tahminleme: Prophet ile "
      "fiyat öngörüsü ve yapısal kırılma tespiti. Bu sistemin asıl muhatabı "
      "Rekabet Kurumu, Ticaret Bakanlığı ve Merkez Bankası gibi "
      "düzenleyiciler; ama şeffaflık sağladığı için nihayetinde tüketiciye "
      "de hizmet ediyor.")
footer(s)

# ===========================================================================
# SLAYT 5 — LITERATURE
# ===========================================================================
s = new_slide()
title_bar(s, "Literature & Existing Solutions", "How others approach this")
left = [
    ("Academic foundation", 2),
    ("Peltzman (2000) — \"Prices Rise Faster Than They Fall\": asymmetric "
     "transmission is pervasive across markets.", 0),
    ("Bacon (1991) — coined \"rockets and feathers\" for fuel retail "
     "pricing; the canonical reference for our asymmetry model.", 0),
    ("Meyer & von Cramon-Taubadel (2004) — survey of asymmetric price "
     "transmission methods (we use the Error-Correction approach).", 0),
]
right = [
    ("Comparable products & data", 2),
    ("TÜİK / TCMB publish CPI food indices — aggregate only, no hal↔retail "
     "linkage and no province-level margin.", 0),
    ("Commercial price-tracker apps (Marketfiyatı, Cimri) compare retail "
     "shelves — but never expose the wholesale margin.", 0),
    ("Remote-sensing crop monitors use Sentinel-2 NDVI to gauge crop "
     "health — an input signal we plan to fuse in (roadmap).", 0),
]
bullets(s, left, 0.7, 1.6, 6.0, 5.0, size=16, gap=11)
bullets(s, right, 6.95, 1.6, 5.7, 5.0, size=16, gap=11)
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7),
                          Inches(5.95), Inches(11.93), Inches(0.8))
_solid(band, NAVY)
band.shadow.inherit = False
tfb = band.text_frame
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfb, "Our gap: no existing product links wholesale & retail to measure "
          "the margin at province scale. That is GıdaRadar's contribution.",
     15, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "5 — Literatür",
      "Hocamızın geçen sunumda haklı olarak istediği gibi, konumuzu "
      "literatürde nasıl çalışıldığına baktık. Akademik temelimiz net: "
      "Peltzman'ın 2000 tarihli, fiyatlar düşerken yükseldiğinden daha yavaş "
      "iner tezi, asimetrik geçişin neredeyse her piyasada görüldüğünü "
      "gösteriyor. Bacon 1991'de rockets and feathers terimini akaryakıt "
      "fiyatları için ortaya atan isim; bizim asimetri modelimizin çıkış "
      "referansı. Meyer ve von Cramon-Taubadel ise asimetrik fiyat geçişi "
      "yöntemlerinin derlemesini yapmış; biz onların önerdiği hata düzeltme "
      "modeli yaklaşımını kullanıyoruz. Mevcut çözümlere bakınca: TÜİK ve "
      "TCMB gıda enflasyon endeksi yayınlıyor ama sadece toplulaştırılmış, "
      "hal-perakende bağlantısı yok. Cimri, Marketfiyatı gibi ticari "
      "uygulamalar raf fiyatlarını karşılaştırıyor ama toptan marjı asla "
      "göstermiyor. Sentinel-2 NDVI ile bitki sağlığı izleyen uydu "
      "çözümleri var; bunu ileride girdi sinyali olarak eklemeyi "
      "planlıyoruz. Özetle: toptan ve perakendeyi birleştirip il ölçeğinde "
      "marjı ölçen mevcut bir ürün yok. GıdaRadar'ın katkısı tam burada.")
footer(s)

# ===========================================================================
# SLAYT 6 — DATA SOURCES
# ===========================================================================
s = new_slide()
title_bar(s, "Data Sources", "9 independent feeds")
rows = [
    ("Source", "Content", "Format", "Cadence"),
    ("marketfiyati.org.tr", "Retail prices · 6 chains · 81 cities", "REST JSON", "Daily"),
    ("İBB & Harman Hal", "Wholesale (hal) prices · 10 real cities", "HTML / CSV", "Daily"),
    ("GDELT (BigQuery)", "Global food/agri news tone & themes", "Parquet", "15-min batch"),
    ("EPİAŞ", "Electricity & gas market · 26 datasets", "REST JSON", "Hourly"),
    ("TCMB EVDS", "FX rates, CPI-food, credit rates", "REST JSON", "Daily / Monthly"),
    ("Open-Meteo", "Weather · 81 cities · frost/heat/rain", "REST JSON", "Hourly"),
    ("Akaryakıt", "Fuel prices · 81 cities (logistics cost)", "REST JSON", "Daily"),
    ("Yahoo Finance", "Global commodities (wheat, brent...)", "REST JSON", "Daily"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(1.5),
                         Inches(11.93), Inches(4.5)).table
widths = [3.0, 4.9, 2.1, 1.93]
for i, w in enumerate(widths):
    tbl.columns[i].width = Inches(w)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(15 if r else 14)
        p.runs[0].font.name = "Calibri"
        if r == 0:
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            p.runs[0].font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r % 2 else WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
tf = textbox(s, 0.7, 6.15, 11.93, 0.8)
para(tf, "5 formats · 5 cadences (15-min → monthly) · merged on a common "
         "(date, city, product) key in the Silver layer.", 16, GREY,
     italic=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "6 — Veri Kaynakları",
      "GıdaRadar dokuz bağımsız veri kaynağını birleştiriyor. İki ana "
      "fiyat kaynağımız var: marketfiyatı.org.tr'den 6 zincir, 81 ilde "
      "perakende fiyatları; İBB ve Harman'dan 10 gerçek ilde hal yani "
      "toptan fiyatları. Bunların etrafına bağlam veren kaynaklar "
      "ekliyoruz: GDELT'ten gıda ve tarım haberlerinin tonu, EPİAŞ'tan "
      "26 veri setiyle elektrik ve gaz piyasası, TCMB EVDS'den döviz kuru "
      "ve gıda enflasyonu, Open-Meteo'dan 81 ilin hava verisi, 81 ilin "
      "akaryakıt fiyatları ve Yahoo Finance'ten buğday, brent gibi küresel "
      "emtialar. İşin zorluğu sayıda değil çeşitlilikte: 5 farklı format ve "
      "15 dakikadan aylığa kadar 5 farklı güncellenme sıklığı var. Hepsini "
      "Silver katmanında ortak bir tarih-şehir-ürün anahtarında "
      "birleştiriyoruz. Bu heterojenliği yönetmek projenin en kritik "
      "mühendislik işiydi.")
footer(s)

# ===========================================================================
# SLAYT 7 — SCRAPING
# ===========================================================================
s = new_slide()
title_bar(s, "Data Collection — Scraping", "Getting the raw data in")
para(textbox(s, 0.7, 1.4, 12.0, 0.6),
     "Each source needed its own technique — no two feeds were alike.",
     16, GREY, italic=True, first=True)
techs = [
    ("Adaptive grid search", ORANGE,
     "marketfiyati API returns depots near one point only. We tile each "
     "district's bounding box into a grid → depots per district jump from "
     "~14 to 629."),
    ("Headless browser", NAVY,
     "İBB hal portal renders prices with JavaScript → Selenium + headless "
     "Chrome drives the page and reads the table."),
    ("Cloudflare bypass", TEAL,
     "Harman hal is behind Cloudflare → curl_cffi with Chrome TLS "
     "fingerprint impersonation gets through cleanly."),
    ("BigQuery extract", GOLD,
     "GDELT GKG is queried in BigQuery (themes × countries) and written "
     "straight to S3 Parquet — parallel workers, 1 TB/mo free tier."),
]
x, y = 0.7, 2.15
for i, (t, c, d) in enumerate(techs):
    cx = x + (i % 2) * 6.1
    cy = y + (i // 2) * 1.95
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx),
                              Inches(cy), Inches(5.85), Inches(1.75))
    _solid(card, LIGHT)
    card.shadow.inherit = False
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx),
                               Inches(cy), Inches(0.16), Inches(1.75))
    _solid(strip, c)
    strip.shadow.inherit = False
    tf = textbox(s, cx + 0.32, cy + 0.13, 5.35, 1.55)
    para(tf, t, 17, c, bold=True, first=True, space_after=5)
    para(tf, d, 13.5, DARK)
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7),
                          Inches(6.05), Inches(11.93), Inches(0.78))
_solid(band, NAVY)
band.shadow.inherit = False
tfb = band.text_frame
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfb, "Shared backbone: async I/O · exponential-backoff rate limiting · "
          "incremental state.json (only scrape what is stale).", 14, WHITE,
     bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "7 — Veri Toplama / Scraping",
      "Hocamızın bu sunumda görmek istediği bir aşama da verinin nasıl "
      "toplandığıydı. Burada şunu vurgulamak isteriz: her kaynak kendi "
      "tekniğini gerektirdi, hiçbiri birbirine benzemiyordu. Dört örnek. "
      "Birincisi adaptif grid search: marketfiyatı API'si sadece tek bir "
      "koordinata yakın depoları döndürüyor, ilçenin uzak mahalleleri "
      "eksik kalıyordu. İlçenin sınır kutusunu bir ızgaraya bölüp her "
      "noktadan sorgu attık; bir ilçedeki depot sayısı 14'ten 629'a çıktı. "
      "İkincisi headless tarayıcı: İBB hal portalı fiyatları JavaScript "
      "ile çiziyor, bu yüzden Selenium ve headless Chrome ile sayfayı "
      "sürüp tabloyu okuduk. Üçüncüsü Cloudflare aşımı: Harman hal "
      "Cloudflare arkasında, curl_cffi ile Chrome'un TLS parmak izini "
      "taklit ederek sorunsuz geçtik. Dördüncüsü BigQuery: GDELT haber "
      "verisini BigQuery'de tema ve ülke filtreleriyle sorgulayıp doğrudan "
      "S3 Parquet'e yazdık. Hepsinin ortak omurgası ise async I/O, üstel "
      "geri çekilmeli rate limiting ve incremental state dosyası — yani "
      "sadece bayatlamış veriyi yeniden çekiyoruz, her seferinde her şeyi "
      "değil.")
footer(s)

# ===========================================================================
# SLAYT 8 — SYNTHETIC DATA
# ===========================================================================
s = new_slide()
title_bar(s, "Synthetic Backfill", "Reaching 2016 with rigor")
para(textbox(s, 0.7, 1.4, 12.2, 0.95),
     "Retail history before 2026 is not public, and hal data is real for "
     "only 10 cities. To study the pandemic and long trends we generate a "
     "rigorous 2016–2026 synthetic backfill for the rest.",
     16, DARK, first=True)
# formul akisi
formula = [("Base price\n(2026 real scrape)", TEAL),
           ("× Inflation deflator\n(TCMB fresh-produce CPI)", NAVY),
           ("× Seasonal profile\n(per-crop harvest curve)", ORANGE),
           ("× Deterministic noise\n(seeded ±4-5%)", GOLD)]
x = 0.7
for i, (t, c) in enumerate(formula):
    box = chip(s, x, 2.6, 2.7, 1.15, t, c, size=13)
    x += 2.78
    if i < 3:
        arrow(s, x - 0.16, 3.02, 0.28, color=GREY)
rows = [
    ("Inflation anchor", "TCMB fresh fruit & vegetable index — 26.3× "
     "cumulative food inflation from 2016 to 2026."),
    ("Reproducible", "Noise seed = hash(product+date+depot) → identical "
     "output on every run, fully auditable."),
    ("Honestly labelled", "Every synthetic row carries veri_turu = "
     "'sentetik'; real vs. synthetic is never mixed silently."),
]
bullets(s, [(t, 2) if False else (t + " — " + d, 0) for t, d in rows],
        0.7, 4.15, 12.0, 2.0, size=16, gap=12)
notes(s, "8 — Sentetik Veri",
      "Bir veri sorunumuz vardı: perakende fiyat geçmişi 2026 öncesi kamuya "
      "açık değil ve hal verisi yalnızca 10 ilde gerçek. Pandemiyi ve uzun "
      "dönem trendleri inceleyebilmek için kalan iller ve geçmiş yıllar "
      "adına titiz bir 2016-2026 sentetik backfill üretiyoruz. Formül dört "
      "çarpandan oluşuyor: 2026'daki gerçek scrape'ten gelen taban fiyat, "
      "TCMB taze meyve-sebze endeksinden türetilen enflasyon deflatörü, her "
      "ürünün hasat eğrisine göre mevsimsel profil ve son olarak "
      "deterministik gürültü. Üç nokta önemli. Enflasyon çıpamız gerçek: "
      "TCMB endeksine göre 2016'dan 2026'ya gıda enflasyonu 26 katı aşmış, "
      "bunu birebir kullanıyoruz. İkincisi tekrarlanabilir: gürültünün "
      "tohumu ürün, tarih ve depot'tan hash ile üretiliyor, yani script her "
      "çalıştığında aynı çıktıyı veriyor, denetlenebilir. Üçüncüsü dürüst "
      "etiketleme: her sentetik satır veri_turu alanında sentetik olarak "
      "işaretli; gerçek ve sentetik veri asla sessizce karışmıyor. Hocamız "
      "da değinmişti, herkes sentetik veri üretiyor; biz de üretiyoruz ama "
      "kaynağı gerçek enflasyona bağlı ve şeffaf şekilde işaretli.")
footer(s)

# ===========================================================================
# SLAYT 9 — BIG DATA SCALE
# ===========================================================================
s = new_slide()
title_bar(s, "Big Data Scale", "Built to grow")
stats = [("81", "provinces covered"),
         ("9", "independent sources"),
         ("2016–2026", "time span"),
         ("~283 M", "rows at full backfill")]
x = 0.7
for v, lbl in stats:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                              Inches(1.6), Inches(2.86), Inches(1.9))
    _solid(card, NAVY)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, v, 32, GOLD, bold=True, first=True, align=PP_ALIGN.CENTER,
         space_after=2)
    para(tf, lbl, 13, WHITE, align=PP_ALIGN.CENTER)
    x += 3.04
body = [
    "Today the working set is tens of GB — but the design target is the "
    "500 GB+ regime where single-machine processing breaks down.",
    "The full 2016–2026 market backfill alone is ~283 M rows / ~52 GB of "
    "Parquet; with all sources it scales straight into true Big Data.",
    "Same code path runs a 1-year demo subset or the full backfill — only "
    "the --start-date flag changes. Volume is a parameter, not a rewrite.",
]
bullets(s, body, 0.7, 3.95, 12.0, 2.7, size=18, gap=14)
notes(s, "9 — Büyük Veri Ölçeği",
      "Projenin ölçeğini birkaç rakamla özetleyelim: 81 il, 9 bağımsız "
      "kaynak, 2016'dan 2026'ya uzanan zaman aralığı ve tam backfill'de "
      "yaklaşık 283 milyon satır. Şunu dürüstçe söyleyelim: bugün aktif "
      "çalıştığımız veri kümesi onlarca gigabayt seviyesinde. Ama tasarım "
      "hedefimiz 500 gigabayt ve üzerindeki bölge — yani tek makinede "
      "işlemenin çöktüğü ölçek. Yalnızca market backfill'i bile 283 milyon "
      "satır, 52 gigabayt Parquet; tüm kaynaklar eklendiğinde doğrudan "
      "gerçek büyük veri rejimine giriyoruz. En önemlisi: aynı kod hem 1 "
      "yıllık demo alt kümesini hem de tam backfill'i çalıştırıyor, sadece "
      "start-date parametresi değişiyor. Yani hacim bizim için bir yeniden "
      "yazım değil, sadece bir parametre. Mimariyi en baştan bunun için "
      "kurduk.")
footer(s)

# ===========================================================================
# SLAYT 10 — END-TO-END DATA PIPELINE (diyagram, tek sayfada mimari)
# ===========================================================================
s = new_slide()
title_bar(s, "End-to-End Data Pipeline", "Medallion on AWS · from source to dashboard")
# pipeline_diagram.png — 1939x811, aspect ~2.39:1
picture(s, "pipeline_diagram.png", 0.4, 1.4, w=12.5)
caption(s,
        "9 sources  →  Apache NiFi + Kafka  →  S3 Bronze / Silver / Gold "
        "(Parquet, year/month)  →  Elasticsearch + Kibana",
        0.5, 6.75, 12.5)
notes(s, "Veri Pipeline'ı (Uçtan Uca)",
      "Bu tek diyagram, projedeki tüm bileşenlerin nasıl bir araya "
      "geldiğini özetliyor — uçtan uca medallion mimarisi, tamamı AWS "
      "üzerinde. Soldan başlayalım. 9 bağımsız veri kaynağımız var: "
      "marketfiyati.org.tr (perakende, REST async), İBB İstanbul Hali ve "
      "Harman Hali (toptan; Selenium ve curl_cffi), TCMB EVDS (kur ve "
      "enflasyon), EPİAŞ Transparency (26 elektrik veri seti), Open-Meteo "
      "(saatlik hava), GDELT (15 dakikalık haber), commodities ve "
      "akaryakıt. Orkestrasyon tarafında — hocamızın özellikle istediği "
      "şekilde Airflow yerine Apache NiFi kullandık; NiFi ingestion'ı "
      "sürekli batch'ler halinde sürüyor, Kafka ise dayanıklı olay "
      "günlüğü görevini üstleniyor. Veri buradan AWS S3'e Parquet "
      "formatında Bronze katmanına iniyor — ham, ama sıkıştırılmış ve "
      "kolonlanmış halde; yıl-ay bazlı Hive partition'larıyla. PySpark "
      "joblarımız Bronze'dan Silver'a entity resolution (hal ürün adı "
      "ile market slug'ını fuzzy match ile eşleme), birim "
      "standartlaştırma (her şey kilograma) ve mekansal eşleme yapıyor. "
      "Silver'dan Gold'a ise 8 analiz tablosu üretiyoruz: günlük marj, "
      "rockets and feathers asimetrisi, hava şok yayılımı, market ve "
      "hal için şehirler arası fiyat farkı, makro korelasyon, haber "
      "korelasyonu ve Prophet tabanlı fiyat tahmini. Tüm hesaplama "
      "PySpark 3.5.1 ile yapılıyor — geliştirme için EC2 local mode, "
      "ağır 10 yıllık koşular için AWS EMR transient cluster: işi yap, "
      "kapan, faturayı sıfırla. Sunum tarafında — yine hocamızın "
      "yönlendirmesiyle Athena yerine Elasticsearch + Kibana tercih "
      "ettik; Gold tabloları ES'e indeksleniyor, Kibana dashboard'ları "
      "interaktif keşif sağlıyor. Diyagramın alt bandı mimarinin altı "
      "özelliğini hatırlatıyor: resilient ve scalable, medallion, "
      "partitioned lake, PySpark powered, curated analytics, search ve "
      "visualize. Yani 9 farklı kaynak ve format, tek bir analitik "
      "katmana dönüşüp interaktif görselleştirmeye ulaşıyor.")
footer(s)

# ===========================================================================
# SLAYT 12 — BRONZE & PARQUET COMPRESSION
# ===========================================================================
s = new_slide()
title_bar(s, "Bronze Layer & Parquet Compression", "Why Parquet")
picture(s, "compression.png", 0.5, 1.5, w=8.0)
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7),
                           Inches(1.5), Inches(4.1), Inches(4.05))
_solid(panel, LIGHT)
panel.shadow.inherit = False
tf = textbox(s, 8.95, 1.7, 3.65, 3.7)
para(tf, "WHY IT MATTERS", 14, TEAL, bold=True, first=True, space_after=9)
for t in ["Columnar storage — read only the columns a query needs.",
          "Type-aware compression — dictionary + run-length encoding.",
          "Partition pruning — WHERE year=2025 skips whole folders.",
          "hal_all: 1028 MB CSV → 62 MB Parquet ≈ 17× smaller.",
          "Smaller files → shorter EMR runtime → lower AWS cost."]:
    para(tf, t, 14, DARK, bullet="▸", space_after=8)
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5),
                          Inches(5.8), Inches(12.3), Inches(0.95))
_solid(band, NAVY)
band.shadow.inherit = False
tfb = band.text_frame
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfb, "Bronze keeps data raw — but as Parquet. Conversion shrinks "
          "storage up to ~17× and cuts every downstream compute bill.",
     15, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "11 — Bronze ve Parquet Sıkışması",
      "Bronze katmanı veriyi ham tutar — ama formatı Parquet. Hocamız "
      "özellikle bunu görmek istemişti: veri Bronze'a girince ne kadar "
      "küçülüyor? Soldaki grafikte kaynak kaynak ham boyut ile Parquet "
      "boyutunu karşılaştırıyoruz. En çarpıcı örnek hal verisi: 1028 "
      "megabayt CSV, Parquet'e dönünce 62 megabayta iniyor — yaklaşık 17 "
      "kat küçülme. Akaryakıtta 12 kat, EPİAŞ'ta 6 kat sıkışma var. Bu "
      "sihir değil; Parquet sütun bazlı saklıyor, yani bir sorgu sadece "
      "ihtiyacı olan sütunları okuyor. Tip-farkında sıkıştırma, sözlük ve "
      "run-length kodlama tekrar eden değerleri katlıyor. Ayrıca "
      "partition pruning sayesinde yıl-2025 filtreli bir sorgu diğer "
      "klasörleri hiç açmıyor. Bunun pratik sonucu şu: daha küçük dosya, "
      "daha kısa EMR çalışma süresi, daha düşük AWS faturası. Yani format "
      "seçimi doğrudan maliyet kararı.")
footer(s)

# ===========================================================================
# SLAYT 12 — SILVER & ENTITY RESOLUTION
# ===========================================================================
s = new_slide()
title_bar(s, "Silver Layer & Entity Resolution", "The hardest problem")
para(textbox(s, 0.7, 1.35, 12.2, 0.6),
     "Hal and retail systems share no product IDs. \"Domates Sofralık Sera\" "
     "must be matched to \"salkim-domates-1-kg\".", 16, DARK, bold=True,
     first=True)
# eslestirme akisi
flow = [("Hal name +\nRetail title", GREY),
        ("Pre-filter\ncandidate pairs", GOLD),
        ("LLM judge\n(Claude Haiku)", TEAL),
        ("Human review\n→ mapping CSV", NAVY)]
x = 0.7
for i, (t, c) in enumerate(flow):
    chip(s, x, 2.15, 2.6, 1.0, t, c, size=12.5)
    x += 2.72
    if i < 3:
        arrow(s, x - 0.18, 2.5, 0.28, color=GREY)
cols = [
    ("Entity resolution", 2),
    ("LLM-based matching (Claude Haiku) — modern language model instead of "
     "classic TF-IDF, per the professor's guidance.", 0),
    ("Batched & prompt-cached → fast, cheap, explainable; human review "
     "keeps the final mapping CSV trustworthy.", 0),
    ("Unit & alignment", 2),
    ("Unit standardisation — Kasa / Bağ / Adet / Gram all normalised to a "
     "1 kg basis via conversion UDFs.", 0),
    ("Temporal & spatial alignment — as-of joins across 15-min → monthly "
     "feeds; retail joined to its nearest hal.", 0),
]
bullets(s, cols, 0.7, 3.5, 12.1, 3.1, size=15, gap=10)
notes(s, "12 — Silver ve Entity Resolution",
      "Silver katmanı projenin en zor mühendislik problemini barındırıyor: "
      "entity resolution. Hal sistemiyle perakende sisteminin ortak ürün "
      "kimliği yok. Haldeki Domates Sofralık Sera ile marketteki "
      "salkim-domates-1-kg aynı ürün ama isimleri tamamen farklı; bunları "
      "eşleştirmek gerekiyor. Akış şöyle: önce hal adı ile market başlığını "
      "alıyoruz, bir ön filtreyle aday çiftleri çıkarıyoruz, sonra bir dil "
      "modeli yani Claude Haiku bu çiftleri yargılıyor, en sonda insan "
      "incelemesiyle nihai eşleştirme CSV'si oluşuyor. Burada hocamızın "
      "yönlendirmesi belirleyici oldu: klasik TF-IDF yerine modern bir dil "
      "modeli kullandık. Çağrılar batch halinde ve prompt-cache'li, bu yüzden "
      "hem hızlı hem ucuz hem de açıklanabilir; insan incelemesi de nihai "
      "eşleştirmeyi güvenilir kılıyor. Silver'da iki iş daha var: birim "
      "standardizasyonu — kasa, bağ, adet, gram, hepsini bir kilogram "
      "tabanına çeviren UDF'ler — ve zamansal-mekânsal hizalama: 15 "
      "dakikalıktan aylığa kadar farklı sıklıktaki kaynakları as-of join'le "
      "birleştirip her perakende noktasını en yakın hale bağlıyoruz.")
footer(s)

# ===========================================================================
# SLAYT 13 — EMR TRANSIENT
# ===========================================================================
s = new_slide()
title_bar(s, "EMR Transient Cluster", "Spin up, process, shut down")
quote = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7),
                           Inches(1.45), Inches(11.93), Inches(0.95))
_solid(quote, TEAL)
quote.shadow.inherit = False
tfq = quote.text_frame
tfq.vertical_anchor = MSO_ANCHOR.MIDDLE
tfq.word_wrap = True
para(tfq, "\"I open the cluster, process the data, and shut it down\" — "
          "compute only exists while a job runs. Zero idle cost.",
     17, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
flow = [("preflight.py\nstatic checks", GREY),
        ("launch_demo.sh\ncreate cluster", GOLD),
        ("11 Spark steps\nBronze→Silver→Gold", TEAL),
        ("--auto-terminate\ncluster dies", ORANGE)]
x = 0.7
for i, (t, c) in enumerate(flow):
    chip(s, x, 2.75, 2.7, 1.05, t, c, size=12.5)
    x += 2.82
    if i < 3:
        arrow(s, x - 0.2, 3.12, 0.3, color=GREY)
rows = [
    "preflight.py validates script syntax, deps.zip and S3 paths BEFORE "
    "any cluster starts — fail early, pay nothing.",
    "Cluster: emr-7.2.0 · Spark 3.5.1 · 1 master + 3 core nodes · the "
    "smoke-test step terminates the cluster instantly if it fails.",
    "A full Bronze→Silver→Gold run costs roughly $1 and auto-terminates — "
    "no forgotten cluster billing overnight.",
]
bullets(s, rows, 0.7, 4.15, 12.1, 2.5, size=16, gap=13)
notes(s, "13 — EMR Transient Cluster",
      "Hesaplama tarafında stratejimiz transient yani geçici EMR cluster. "
      "Mantığı tek cümleyle şu: cluster'ı açıyorum, veriyi işliyorum, "
      "kapatıyorum. Yani işlem gücü yalnızca bir iş koşarken var oluyor, "
      "boşta hiç maliyet yok. Akış dört adım. Önce preflight script'i: "
      "cluster daha açılmadan script sözdizimini, deps.zip paketini ve S3 "
      "yollarını statik olarak doğruluyor — erken hata ver, hiç para ödeme "
      "mantığı. Sonra launch_demo script'i cluster'ı yaratıyor. Cluster "
      "emr-7.2.0, Spark 3.5.1, bir master ve üç core node. Sonra 11 Spark "
      "adımı Bronze'dan Silver'a, Silver'dan Gold'a çalışıyor; ilk adım bir "
      "smoke testi ve başarısız olursa cluster anında kendini sonlandırıyor. "
      "En sonda auto-terminate ile cluster ölüyor. Tam bir Bronze-Silver-Gold "
      "koşusu yaklaşık 1 dolara mal oluyor ve kendi kendine kapanıyor; gece "
      "boyu unutulup fatura kabartan bir cluster riski yok.")
footer(s)

# ===========================================================================
# SLAYT 14 — EMR RUN PERFORMANCE
# ===========================================================================
s = new_slide()
title_bar(s, "EMR Run Performance", "The 1-year pipeline run")
para(textbox(s, 0.7, 1.35, 12.2, 0.6),
     "We ran the full Bronze→Silver→Gold pipeline on a 1-year subset "
     "(2025-05-20 → 2026-05-20) on AWS EMR.", 16, DARK, bold=True, first=True)
stats = [("emr-7.2.0", "Spark 3.5.1 · Py 3.9"),
         ("1 + 3", "master + core nodes"),
         ("11", "Spark steps"),
         ("~30–45 min", "end-to-end runtime"),
         ("≈ $1", "total cost / run")]
x = 0.7
for v, lbl in stats:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                              Inches(2.05), Inches(2.28), Inches(1.5))
    _solid(card, NAVY)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, v, 21, GOLD, bold=True, first=True, align=PP_ALIGN.CENTER,
         space_after=2)
    para(tf, lbl, 11.5, WHITE, align=PP_ALIGN.CENTER)
    x += 2.42
rows = [
    "Step 0 is a smoke test — processes a single day without writing; if "
    "it fails the cluster terminates, so the other 10 steps cost nothing.",
    "Steps 1–10 chain market_silver → silver_joined → gdelt_silver → the "
    "7 Gold analyses, with ActionOnFailure=CONTINUE so one failure does "
    "not abort the rest.",
    "Result: the entire medallion pipeline — millions of rows across 9 "
    "sources — completes in well under an hour for about one dollar.",
]
bullets(s, rows, 0.7, 3.85, 12.1, 2.6, size=15.5, gap=12)
notes(s, "14 — EMR Çalıştırma Performansı",
      "Bu pipeline'ı sadece tasarlamadık, gerçekten çalıştırdık. Tam "
      "Bronze-Silver-Gold hattını 1 yıllık bir alt küme üzerinde — yani "
      "2025 Mayıs'tan 2026 Mayıs'a — AWS EMR'da koşturduk. Rakamlar: cluster "
      "emr-7.2.0, Spark 3.5.1, Python 3.9. Bir master artı üç core node. "
      "Toplam 11 Spark adımı. Uçtan uca çalışma süresi 30 ila 45 dakika. Ve "
      "koşu başına toplam maliyet yaklaşık 1 dolar. Akışın mantığı şöyle: "
      "sıfırıncı adım bir smoke testi, tek bir günü hiçbir şey yazmadan "
      "işliyor; başarısız olursa cluster kapanıyor, böylece diğer 10 adım "
      "hiç para harcamıyor. Birden ona kadarki adımlar market_silver, "
      "silver_joined, gdelt_silver ve ardından 7 Gold analizini zincirleme "
      "çalıştırıyor; ActionOnFailure CONTINUE olduğu için bir adım hata "
      "verse bile diğerleri devam ediyor. Sonuç şu: 9 kaynaktan gelen "
      "milyonlarca satırı kapsayan tüm medallion pipeline, bir saatin epey "
      "altında ve yaklaşık bir dolara tamamlanıyor. Hocamızın istediği "
      "büyük veri ile baş etme yetkinliğini somut olarak gösteren kısım bu.")
footer(s)

# ===========================================================================
# SLAYT 15 — UNDER THE HOOD: SPARK
# ===========================================================================
s = new_slide()
title_bar(s, "Under the Hood — Spark Workload", "What runs in the background")
para(textbox(s, 0.7, 1.35, 12.2, 0.6),
     "Behind each step is a distributed Spark job — here is what they "
     "actually do.", 16, GREY, italic=True, first=True)
items = [
    ("Trunk join", ORANGE,
     "silver_joined performs a FULL OUTER JOIN of retail × hal on "
     "(date, city, product); the small mapping CSV is broadcast to every "
     "executor to avoid a shuffle."),
    ("Window functions", NAVY,
     "daily_margin computes 7-day rolling averages with Spark Window "
     "functions — partitioned by city/product, no full shuffle."),
    ("Distributed model fit", TEAL,
     "rockets_feathers, shock_propagation and Prophet fit one statistical "
     "model per product group via applyInPandas — each executor fits in "
     "parallel, not the driver."),
    ("Spark optimisations", GOLD,
     "Adaptive Query Execution coalesces shuffle partitions; year/month "
     "partition pruning skips S3 folders the query never needs."),
]
x, y = 0.7, 2.05
for i, (t, c, d) in enumerate(items):
    cx = x + (i % 2) * 6.1
    cy = y + (i // 2) * 2.0
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx),
                              Inches(cy), Inches(5.85), Inches(1.8))
    _solid(card, LIGHT)
    card.shadow.inherit = False
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx),
                               Inches(cy), Inches(0.16), Inches(1.8))
    _solid(strip, c)
    strip.shadow.inherit = False
    tf = textbox(s, cx + 0.32, cy + 0.12, 5.35, 1.6)
    para(tf, t, 16, c, bold=True, first=True, space_after=4)
    para(tf, d, 12.5, DARK)
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7),
                          Inches(6.15), Inches(11.93), Inches(0.68))
_solid(band, NAVY)
band.shadow.inherit = False
tfb = band.text_frame
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfb, "Heavy lifting is distributed across the cluster — the driver "
          "only collects small, already-aggregated results.", 14, WHITE,
     bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "15 — Arka Plandaki Spark İşleri",
      "Her EMR adımının arkasında dağıtık bir Spark işi var; kısaca bu "
      "işlerin gerçekte ne yaptığını anlatalım. Birincisi gövde join'i: "
      "silver_joined, perakende ile hal verisini tarih-şehir-ürün anahtarı "
      "üzerinden FULL OUTER JOIN yapıyor; küçük eşleştirme CSV'si her "
      "executor'a broadcast ediliyor, böylece pahalı bir shuffle'dan "
      "kaçınıyoruz. İkincisi pencere fonksiyonları: daily_margin, 7 günlük "
      "hareketli ortalamaları Spark Window fonksiyonlarıyla hesaplıyor, "
      "şehir ve ürüne göre bölümlenmiş şekilde, tam shuffle olmadan. "
      "Üçüncüsü dağıtık model fit'i: rockets_feathers, shock_propagation ve "
      "Prophet, her ürün grubu için ayrı bir istatistiksel model kuruyor ve "
      "bunu applyInPandas ile yapıyor — yani modeller driver'da değil, her "
      "executor'da paralel olarak fit ediliyor. Dördüncüsü Spark "
      "optimizasyonları: Adaptive Query Execution shuffle partition'larını "
      "birleştiriyor, yıl-ay partition pruning ise sorgunun ihtiyaç "
      "duymadığı S3 klasörlerini hiç açmıyor. Özet: ağır iş cluster'a "
      "dağıtılıyor, driver yalnızca küçük ve önceden toplulaştırılmış "
      "sonuçları topluyor.")
footer(s)

# ===========================================================================
# SLAYT 16 — METHODOLOGY
# ===========================================================================
s = new_slide()
title_bar(s, "Methodology", "The analytical toolbox")
rows = [
    ("Method", "Used for", "What it does"),
    ("Asymmetric Error-\nCorrection Model", "Rockets & Feathers",
     "Engle-Granger ECM splitting up/down moves → asymmetry score = "
     "Σβ⁺ / |Σβ⁻|."),
    ("Event study +\nlag detection", "Shock propagation",
     "Detect frost/heat/rain events, count days until price crosses "
     "+10% over baseline."),
    ("Coefficient of\nvariation & spread", "Price inequality",
     "Std/mean and (max−min)/avg of one product across 81 provinces."),
    ("Pearson lag\ncorrelation", "Macro & news drivers",
     "Correlate price change with macro/news series at lags 0–30 days; "
     "report best lag + p-value."),
    ("Prophet additive\nmodel", "Forecasting",
     "Trend + yearly seasonality + changepoints → 30-day forecast with "
     "uncertainty bands."),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.5),
                         Inches(11.93), Inches(5.0)).table
for i, w in enumerate([2.9, 2.7, 6.33]):
    tbl.columns[i].width = Inches(w)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        col0 = (c == 0)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(14 if r else 13)
                run.font.name = "Calibri"
                if r == 0:
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                else:
                    run.font.color.rgb = NAVY if col0 else DARK
                    run.font.bold = col0
        cell.fill.solid()
        cell.fill.fore_color.rgb = (
            NAVY if r == 0 else (LIGHT if r % 2 else WHITE))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1)
notes(s, "16 — Metodoloji",
      "Hocamız metodolojinin net görünmesini özellikle istemişti, bu yüzden "
      "kullandığımız analitik araç kutusunu tek tabloda topladık. Beş "
      "yöntem var. Birincisi asimetrik hata düzeltme modeli — Engle-Granger "
      "ECM — Rockets and Feathers için; fiyat hareketlerini yukarı ve aşağı "
      "diye ikiye ayırıp asimetri skorunu pozitif katsayılar toplamının "
      "negatiflerin mutlak değerine oranı olarak hesaplıyor. İkincisi olay "
      "çalışması ve gecikme tespiti — şok yayılımı için; don, sıcak, yağış "
      "olaylarını tespit edip fiyatın baz çizgisinin yüzde 10 üstüne "
      "çıkması kaç gün sürdüğünü sayıyor. Üçüncüsü varyasyon katsayısı ve "
      "spread — fiyat eşitsizliği için; bir ürünün 81 ildeki standart "
      "sapma bölü ortalama ve maksimum eksi minimum bölü ortalama "
      "değerleri. Dördüncüsü Pearson gecikmeli korelasyon — makro ve haber "
      "etkenleri için; fiyat değişimini 0 ila 30 günlük gecikmelerde makro "
      "serilerle ilişkilendirip en iyi gecikmeyi ve p değerini raporluyor. "
      "Beşincisi Prophet additive model — tahminleme için; trend, yıllık "
      "mevsimsellik ve değişim noktalarıyla 30 günlük öngörü ve "
      "belirsizlik bantları üretiyor. Sıradaki slaytlarda bu yöntemlerin "
      "gerçek sonuçlarını göreceğiz.")
footer(s)

# ===========================================================================
# SLAYT 17 — RESULT: DAILY MARGIN
# ===========================================================================
result_slide(
    "Result 1 of 6", "Daily Hal → Retail Margin", "margin_dist.png",
    ["Every daily hal→retail price record in the run is analysed.",
     "Median markup: the shelf price is typically close to double the "
     "wholesale price.",
     "A long right tail — some product/province cases show extreme "
     "margins that pull the mean far above the median.",
     "Average margin varies clearly from province to province."],
    "On the typical product the retailer's price is roughly double what "
    "the wholesaler charged.")
notes(NOTES and prs.slides[-1], "17 — Sonuç: Günlük Marj",
      "İlk sonucumuz günlük hal-perakende marjı. Pipeline'daki tüm günlük "
      "fiyat kayıtlarını analiz ettik — soldaki grafikte medyan turuncu "
      "çizgiyle işaretli. Çıkan tablo şu: tipik bir üründe raf fiyatı, "
      "toptan fiyatının neredeyse iki katı. Dağılımın uzun bir sağ kuyruğu "
      "var; yani bazı ürün ve il kombinasyonlarında marj olağanüstü "
      "açılıyor ve bu uç değerler ortalamayı medyanın epey üstüne çekiyor. "
      "Medyan ile ortalama arasındaki bu fark başlı başına bir bulgu: "
      "dağılım sağa çarpık, marj her yerde aynı değil. Sağdaki grafikte ise "
      "marjın ilden ile belirgin biçimde değiştiğini görüyoruz — en yüksek "
      "marjlı iller, en düşüklerden net biçimde ayrışıyor. Buradan çıkan "
      "içgörü şu: tipik üründe perakendecinin fiyatı, toptancının aldığının "
      "kabaca iki katı. Bu, projenin temel iddiasını rakamla doğruluyor. "
      "Not: grafikteki sayılar canlı veriden geliyor, sunum öncesi "
      "yeniden üretildiğinde güncel değeri yansıtır.")

# ===========================================================================
# SLAYT 18 — RESULT: ROCKETS & FEATHERS
# ===========================================================================
result_slide(
    "Result 2 of 6", "Rockets & Feathers — Asymmetry", "rockets_feathers.png",
    ["Asymmetric ECM fitted per product (≥150 observations each).",
     "Score > 1 = \"rocket\": shelf rises fast on cost increases.",
     "Score < 1 = \"feather\": shelf falls slowly on cost decreases.",
     "Staples like garlic & potato show the strongest rocket behaviour."],
    "Several staple products pass cost increases through far faster than "
    "they pass decreases — measurable consumer-unfavourable asymmetry.")
notes(NOTES and prs.slides[-1], "18 — Sonuç: Rockets & Feathers",
      "İkinci sonucumuz, projenin en ilgi çekici kısmı: Rockets and "
      "Feathers asimetrisi. Her ürün için, en az 150 gözlem şartıyla, "
      "asimetrik hata düzeltme modeli fit ettik. Grafikte dikey kesikli "
      "çizgi 1.0 değerinde, yani simetrik geçişi gösteriyor. Skoru 1'in "
      "üstünde olan ürünler — turuncu çubuklar — roket: maliyet artınca raf "
      "fiyatı hızlı yükseliyor. Skoru 1'in altında olanlar — yeşil çubuklar "
      "— tüy: maliyet düşünce raf fiyatı yavaş iniyor. Sonuç çarpıcı biçimde "
      "iki kutuplu. Kuru sarımsak, bebe patates, patates gibi temel "
      "ürünlerde belirgin roket davranışı var; bu ürünlerde fiyat "
      "artışları, düşüşlere kıyasla çok daha hızlı rafa yansıyor. Buradan "
      "çıkan içgörü net: birçok temel gıda üründe maliyet artışları, "
      "maliyet düşüşlerinden çok daha hızlı tüketiciye geçiyor — yani "
      "ölçülebilir, tüketici aleyhine bir asimetri var. Bu tam olarak "
      "Bacon'ın rockets and feathers dediği etkinin Türkiye gıda perakende "
      "verisinde sayısal kanıtı.")

# ===========================================================================
# SLAYT 19 — RESULT: SHOCK PROPAGATION
# ===========================================================================
result_slide(
    "Result 3 of 6", "Weather Shock Propagation", "shock_lag.png",
    ["Frost / heat / heavy-rain events detected from weather data.",
     "We measure the lag until wholesale, then retail, prices jump +10%.",
     "Wholesale (hal) reacts first; retail follows with extra delay.",
     "Leafy & fruiting vegetables show the largest peak price jumps."],
    "Weather shocks reach the shelf with a measurable delay — enabling an "
    "early-warning window of several days.")
notes(NOTES and prs.slides[-1], "19 — Sonuç: Şok Yayılımı",
      "Üçüncü sonucumuz hava şoku yayılımı. Hava verisinden don, aşırı "
      "sıcak ve şiddetli yağış olaylarını tespit ediyoruz. Sonra her olay "
      "için şu soruyu soruyoruz: fiyatın baz çizgisinin yüzde 10 üstüne "
      "çıkması kaç gün sürdü? Bunu önce hal yani toptan, sonra market yani "
      "perakende için ölçüyoruz. Soldaki histogramda gecikme günlerinin "
      "dağılımı var; hal fiyatı önce tepki veriyor, perakende ise ek bir "
      "gecikmeyle onu takip ediyor. Sağdaki grafikte ise hava olaylarına en "
      "duyarlı ürünleri görüyoruz — yapraklı sebzeler ve meyveli sebzeler "
      "en büyük zirve fiyat sıçramalarını yaşıyor. Buradan çıkan içgörü "
      "uygulama açısından çok değerli: hava şokları rafa belirli, ölçülebilir "
      "bir gecikmeyle ulaşıyor. Bu gecikme bir erken uyarı penceresi demek "
      "— don olayını gördüğümüz anda, fiyat henüz rafa yansımadan birkaç "
      "günlük bir öngörü süresi kazanıyoruz. Erken uyarı radarı fikri tam "
      "da bu gecikmenin üstüne kuruluyor.")

# ===========================================================================
# SLAYT 20 — RESULT: PRICE INEQUALITY
# ===========================================================================
result_slide(
    "Result 4 of 6", "Inter-City Price Inequality", "price_inequality.png",
    ["Same product, 81 provinces — how wide is the price gap?",
     "Spread % = (max − min) / average shelf price across cities.",
     "Top products show a large, persistent gap between the cheapest and "
     "the priciest province.",
     "This is the 'spatial speculation' signal made measurable."],
    "Identical products carry wide price gaps between provinces — far "
    "beyond what logistics cost alone can explain.")
notes(NOTES and prs.slides[-1], "20 — Sonuç: Fiyat Eşitsizliği",
      "Dördüncü sonucumuz şehirler arası fiyat eşitsizliği. Soru basit: "
      "aynı ürün, 81 il — fiyat farkı ne kadar açılıyor? Bunu spread yüzdesi "
      "ile ölçüyoruz: iller arasındaki maksimum fiyat eksi minimum fiyat, "
      "bölü ortalama raf fiyatı. Grafikte en yüksek eşitsizliğe sahip "
      "ürünleri görüyoruz; üst sıradaki ürünlerde en ucuz il ile en pahalı "
      "il arasında onlarca puanlık bir fiyat farkı var. Yani birebir aynı "
      "ürün, bir ilde diğerinden belirgin biçimde daha pahalı olabiliyor. "
      "Köşedeki kutuda ortalama varyasyon katsayısını da not ettik. Bu "
      "bulgu, problem slaytında bahsettiğimiz mekânsal spekülasyon "
      "sinyalini ölçülebilir hale getiriyor. Aynı ürünler iller arasında "
      "çift haneli, kimi zaman yüzde 50'ye varan fiyat farkları taşıyor. "
      "Bu fark lojistik maliyetiyle kısmen açıklanabilir ama tamamı değil; "
      "geri kalanı lokasyon bazlı fiyatlama gücüne işaret ediyor ve "
      "düzenleyici için tam da bakılması gereken yer burası.")

# ===========================================================================
# SLAYT 21 — RESULT: MACRO / EXTERNAL DRIVERS
# ===========================================================================
result_slide(
    "Result 5 of 6", "External Drivers — Macro & News", "macro_corr.png",
    ["Food prices correlated with fuel, FX, commodities & electricity.",
     "Pearson correlation tested at lags of 0–30 days; best lag reported.",
     "GDELT news tone is also tested as a leading indicator of price moves.",
     "Macro factors give context the price data alone cannot."],
    "Food prices do not move in isolation — fuel, FX and commodity shocks "
    "are measurable leading signals.")
notes(NOTES and prs.slides[-1], "21 — Sonuç: Dış Etkenler",
      "Beşinci sonucumuz dış etkenler — makro ve haber. Buraya kadar marjı "
      "ve fiyatı kendi içinde inceledik; ama gıda fiyatları boşlukta "
      "hareket etmiyor. Bu analizde gıda fiyatlarını akaryakıt, döviz kuru, "
      "küresel emtialar ve elektrik fiyatlarıyla ilişkilendirdik. Yöntem "
      "Pearson korelasyonu; her makro seri için fiyat değişimine karşı 0 "
      "ila 30 günlük gecikmeleri test edip en güçlü ilişkiyi veren "
      "gecikmeyi raporluyoruz. Ayrıca GDELT haber tonunu da fiyat "
      "hareketlerinin öncü göstergesi olarak test ediyoruz. Grafikte hangi "
      "makro etkenlerin gıda fiyatıyla en yüksek mutlak korelasyona sahip "
      "olduğunu sıraladık. Buradan çıkan içgörü şu: gıda fiyatları izole "
      "değil; akaryakıt, kur ve emtia şokları ölçülebilir öncü sinyaller. "
      "Bu da erken uyarı sistemine ayrı bir boyut katıyor — sadece haldeki "
      "fiyatı değil, onun da öncülü olan makro etkenleri izleyebiliyoruz.")

# ===========================================================================
# SLAYT 22 — RESULT: FORECASTING
# ===========================================================================
result_slide(
    "Result 6 of 6", "Forecasting with Prophet", "forecast_domates.png",
    ["Tomato price modelled over the full 2016–2026 history.",
     "Prophet captures strong yearly seasonality + the rising trend.",
     "Changepoints flag structural breaks in the trend.",
     "A 30-day forecast is produced with an uncertainty band."],
    "The pipeline does not just explain the past — it produces a "
    "forward-looking 30-day price outlook.")
notes(NOTES and prs.slides[-1], "22 — Sonuç: Tahminleme",
      "Altıncı ve son sonucumuz tahminleme. Hocamız özellikle domatesin son "
      "10 yıllık serisini ve Prophet ile gelecek tahminini görmek "
      "istemişti. Grafikte domates fiyatını 2016'dan 2026'ya tüm geçmişiyle "
      "modelledik. Meta'nın Prophet modeli iki şeyi çok net yakalıyor: "
      "güçlü yıllık mevsimsellik — her yıl tekrarlayan o dalga — ve "
      "alttaki yükselen trend. Kesikli dikey çizgiler trend değişim "
      "noktalarını, yani yapısal kırılmaları işaretliyor. En sağda ise "
      "turuncu renkte 30 günlük gelecek tahmini var; etrafındaki gölgeli "
      "alan da belirsizlik bandı. Buradan çıkan içgörü şu: pipeline "
      "yalnızca geçmişi açıklamıyor, ileriye dönük 30 günlük bir fiyat "
      "öngörüsü de üretiyor. Mevsimsellik artı trend artı değişim noktası "
      "yaklaşımı, erken uyarı radarının tahmin ayağını oluşturuyor.")

# ===========================================================================
# SLAYT 23 — PANDEMIC GAP
# ===========================================================================
s = new_slide()
title_bar(s, "Pandemic Gap Analysis", "Did margins widen for good?")
picture(s, "pandemic_gap.png", 0.4, 1.4, w=8.4)
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0),
                           Inches(1.4), Inches(4.0), Inches(4.4))
_solid(panel, LIGHT)
panel.shadow.inherit = False
tf = textbox(s, 9.22, 1.55, 3.6, 4.1)
para(tf, "RESEARCH QUESTION", 13, TEAL, bold=True, first=True, space_after=5)
para(tf, "Did the hal→retail margin widen permanently after the 2020 "
         "pandemic, vs reverting to pre-2020?", 12.5, DARK, space_after=10)
para(tf, "METHOD", 13, TEAL, bold=True, space_after=5)
para(tf, "1,446 product×chain cases compared: avg margin in 2019 vs "
         "each year 2021–2024. Formula:", 12.5, DARK, space_after=4)
para(tf, "gap% = (post − baseline) / |baseline| × 100", 12, NAVY,
     bold=True, italic=True, space_after=10)
para(tf, "FINDING", 13, ORANGE, bold=True, space_after=5)
para(tf, "Median gap NARROWED each year (−28% to −85%). Hypothesis NOT "
         "confirmed at population level.", 12.5, DARK, space_after=10)
para(tf, "CAVEAT", 13, GREY, bold=True, space_after=5)
para(tf, "Market data 2019–2025 is synthetic backfill — direction is "
         "indicative, not conclusive.", 12.5, GREY, italic=True)
band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4),
                          Inches(6.55), Inches(12.5), Inches(0.85))
_solid(band, NAVY)
band.shadow.inherit = False
tfb = band.text_frame
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfb, "Counter-intuitive result: median product-chain margin actually "
          "compressed post-pandemic — only specific products (right panel) "
          "show permanent widening.",
     14, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "Pandemi Gap Analizi",
      "Hocamızın altını çizdiği araştırma sorusu: hal-perakende marjı 2020 "
      "pandemisinden sonra kalıcı olarak açıldı mı, yoksa pandemi öncesi "
      "seviyesine geri döndü mü? Yöntem: her ürün ve zincir için 2019 "
      "ortalama marjı baseline, bunu 2021, 2022, 2023 ve 2024'ün her biriyle "
      "ayrı ayrı karşılaştırdık. Formül: post eksi baseline, bölü mutlak "
      "baseline, çarpı yüz. Toplam 1.446 ürün-zincir vakası, 4 yıl üzerinden. "
      "Hesabı EMR transient cluster'da koşturduk — pandemic_gap step'i "
      "demo run'da unutulmuştu, sadece bu adıma özel küçük bir cluster "
      "açıp tamamladık. Sonuç hipotezimize ters çıktı: medyan gap her yıl "
      "negatif — 2021'de eksi yirmi sekiz, 2022'de eksi seksen beş, 2023'te "
      "eksi altmış altı, 2024'te eksi otuz üç. Yani çoğu ürün-zincir "
      "kombinasyonunda marj 2019'a göre genişlemedi, daraldı. Bu beklenmedik "
      "ama açıklanabilir: pandemi sonrası yüksek enflasyon ortamında hal "
      "fiyatları çoğu üründe perakendeye göre daha hızlı yükselmiş "
      "olabilir; yani perakende marjı görece sıkışmış. Sağdaki panelde "
      "kuralın istisnaları var — bazı spesifik ürünlerde gap gerçekten "
      "kalıcı olarak açılmış, bunlar denetim için ilginç vakalar. Önemli "
      "bir çekince: market verisi 2019-2025 arası sentetik backfill, "
      "yani sonuç yön açısından bilgilendirici ama kesin değil. Gerçek "
      "perakende verisi geçmişe genişletildiğinde bu rakamlar netleşecek.")
footer(s)

# ===========================================================================
# SLAYT 24 — DASHBOARD
# ===========================================================================
s = new_slide()
title_bar(s, "Dashboards", "Elasticsearch + Kibana — live screenshots")
# 2 × 3 grid: 6 thumbnail (4 dashboard, Marj ve Sok ikiser view)
dashes = [
    ("Marj — Genel Bakış",       "marj-genel-bakis-1.png"),
    ("Rockets & Feathers",        "rockets-feathers.png"),
    ("Prophet — Fiyat Tahmini",   "prophet.png"),
    ("Marj — Türkiye Haritası",   "marj-genel-bakis-2.png"),
    ("Şok Yayılım — İstatistik",  "sok-yayilim-1.png"),
    ("Şok Yayılım — Harita",      "sok-yayilim-2.png"),
]
# Width 3.5 — tallest aspect (prophet ~1.55) -> h≈2.26 < row budget
W_CELL = 3.5
H_HEADER = 0.32
GAP_X = 0.18
ROW_PITCH = 2.68
# Grid'i ortala: 3*3.5 + 2*0.18 = 10.86, slayt 13.333 → margin 1.24
X0 = 1.24
Y0 = 1.4
for i, (name, img) in enumerate(dashes):
    col = i % 3
    row = i // 3
    x = X0 + col * (W_CELL + GAP_X)
    y = Y0 + row * ROW_PITCH
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                             Inches(W_CELL), Inches(H_HEADER))
    _solid(hdr, NAVY)
    hdr.shadow.inherit = False
    tfh = hdr.text_frame
    tfh.vertical_anchor = MSO_ANCHOR.MIDDLE
    tfh.margin_left = Inches(0.1); tfh.margin_right = Inches(0.1)
    para(tfh, name, 11, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
    # Görsel — başlığın hemen altında, hücre genişliğinde; yükseklik aspect ile
    picture(s, img, x, y + H_HEADER + 0.03, w=W_CELL)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                          Inches(6.75), SW, Inches(0.7))
_solid(band, NAVY)
band.shadow.inherit = False
tfb = band.text_frame
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
tfb.word_wrap = True
para(tfb, "From a national heat map to a single product-city-day in "
          "seconds — Kibana lets regulators, journalists & policy teams "
          "drill the Gold layer interactively.",
     13.5, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, "Dashboard",
      "Analiz sonuçlarını yalnızca tablo olarak bırakmıyoruz; Gold "
      "tabloları Elasticsearch'e indeksleniyor ve Kibana üzerinden "
      "etkileşimli olarak keşfediliyor. Hocamızın yönlendirmesiyle Athena "
      "yerine ELK yığınını seçtik; çünkü statik bir sorgu değil, üstüne "
      "tıklanıp keşif yapılan bir deneyim istiyoruz. Slaytta dört ana "
      "dashboard'umuzdan gerçek ekran görüntüleri var — beşincisi makro "
      "etkenler panosu, slaytta yer kısıtı nedeniyle göstermedik. Birinci "
      "satır soldan sağa: Marj Genel Bakış'ta KPI'lar, zincir bazlı marj "
      "trendi ve il bazlı top-20 tablosu görüyoruz; Rockets and Feathers "
      "panosunda ürün ve zincir bazlı asimetri skorları artı 30 satırlık "
      "detay tablosu; Prophet panosunda ürün bazlı tahminler ve güven "
      "aralıkları. İkinci satırda: Marj Türkiye Haritası — il bazlı renkli "
      "ısı görünümü, koyu kırmızı yüksek marj demek, Marmara ve Ege'de "
      "yoğun; Şok Yayılım istatistik panosu — toplam 912 bin şok olayı, "
      "olay tipi dağılımı ve hal-market gecikmesi; ve Şok Türkiye "
      "Haritası — pikem değişim yüzdesini il bazlı gösteriyor, Karadeniz "
      "ve Akdeniz'de daha çarpıcı. Kibana'nın asıl değeri şu: ulusal "
      "haritadan başlayıp saniyeler içinde tek bir ürün-şehir-gün "
      "kırılımına inebiliyorsunuz. Yani sistem sadece veri üretmiyor, o "
      "veriyi karar verilebilir hale getiriyor.")
footer(s)

# ===========================================================================
# SLAYT 25 — BUSINESS VALUE
# ===========================================================================
s = new_slide()
title_bar(s, "Business & Strategic Value", "Who benefits")
groups = [
    ("Regulators", NAVY,
     "Competition Authority & Ministry of Trade get real-data evidence of "
     "unfair pricing — asymmetry scores and margin outliers become "
     "actionable enforcement signals."),
    ("Consumers & Society", TEAL,
     "Transparency on where the margin actually sits demystifies food "
     "inflation and counters speculative pricing psychology."),
    ("Economic Strategy", ORANGE,
     "Pinpoints provinces with logistics-cost problems and chains that "
     "exploit asymmetric transmission — input for targeted policy."),
]
y = 1.7
for t, c, d in groups:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7),
                              Inches(y), Inches(11.93), Inches(1.5))
    _solid(card, LIGHT)
    card.shadow.inherit = False
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7),
                             Inches(y), Inches(3.0), Inches(1.5))
    _solid(tag, c)
    tag.shadow.inherit = False
    tft = tag.text_frame
    tft.vertical_anchor = MSO_ANCHOR.MIDDLE
    tft.word_wrap = True
    para(tft, t, 18, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
    tf = textbox(s, 4.0, y + 0.18, 8.4, 1.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, d, 15, DARK, first=True)
    y += 1.67
notes(s, "25 — İş ve Stratejik Değer",
      "Peki bu sistem kime, ne değer üretiyor? Üç kitleye bakalım. "
      "Birincisi düzenleyiciler: Rekabet Kurumu ve Ticaret Bakanlığı, "
      "haksız fiyatlamayı artık gerçek veriyle gösterebilir. Asimetri "
      "skorları ve marj aykırı değerleri, soyut şikâyetler değil, "
      "uygulanabilir denetim sinyalleri haline geliyor. Sunumun başında "
      "gördüğümüz File Market vakası tam da bunun kanıtı: düzenleyici tek "
      "bir vakayı manuel denetimle yakaladı — GıdaRadar bu tür yüzlerce "
      "aykırı vakayı otomatik olarak önüne koyabilir. İkincisi tüketici "
      "ve toplum: marjın gerçekte nerede oluştuğunun şeffaf olması gıda "
      "enflasyonunu anlaşılır kılıyor ve spekülatif fiyatlama psikolojisini "
      "kırıyor. Üçüncüsü ekonomik strateji: sistem hangi illerde lojistik "
      "maliyet sorunu olduğunu ve hangi zincirlerin asimetrik geçişi "
      "fırsata çevirdiğini nokta atışı gösteriyor; bu da hedefli politika "
      "için doğrudan girdi. Özetle GıdaRadar bir akademik egzersiz değil, "
      "somut bir karar destek aracı olarak tasarlandı.")
footer(s)

# ===========================================================================
# SLAYT 26 — ROADMAP
# ===========================================================================
s = new_slide()
title_bar(s, "Roadmap & Future Work", "Where we go next")
items = [
    ("Full backfill", "Run the complete 2016–2026 history on EMR — "
     "unlocking the pandemic-gap result at full scale."),
    ("ML models", "Add XGBoost & LSTM price models alongside Prophet for "
     "richer multi-feature forecasting."),
    ("Satellite signal", "Fuse Sentinel-2 NDVI crop-health data to detect "
     "supply shocks before they reach the hal."),
    ("Continuous streaming", "Promote the NiFi + Kafka path to a fully "
     "continuous pipeline beyond batch."),
    ("Route analysis", "Study transport corridors (e.g. Ankara–Konya) to "
     "separate logistics cost from pricing power."),
    ("More cities & products", "Extend real hal coverage beyond 10 cities "
     "and deepen the product catalogue."),
]
x, y = 0.7, 1.65
for i, (t, d) in enumerate(items):
    cx = x + (i % 2) * 6.1
    cy = y + (i // 3) * 0.0 + (i // 2) * 1.62
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx),
                              Inches(cy), Inches(5.85), Inches(1.45))
    _solid(card, LIGHT)
    card.shadow.inherit = False
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.18),
                             Inches(cy + 0.18), Inches(0.55), Inches(0.55))
    _solid(num, TEAL)
    num.shadow.inherit = False
    tfn = num.text_frame
    tfn.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tfn, str(i + 1), 16, WHITE, bold=True, first=True,
         align=PP_ALIGN.CENTER)
    tf = textbox(s, cx + 0.95, cy + 0.13, 4.75, 1.25)
    para(tf, t, 15, NAVY, bold=True, first=True, space_after=3)
    para(tf, d, 12.5, DARK)
notes(s, "26 — Yol Haritası",
      "Proje burada bitmiyor; yol haritamızda altı başlık var. Birincisi "
      "tam backfill: 2016-2026 tüm geçmişi EMR'da koşturmak — bu, pandemi "
      "gap sonucunu tam ölçekte açacak. İkincisi makine öğrenmesi "
      "modelleri: Prophet'in yanına XGBoost ve LSTM ekleyip çok değişkenli, "
      "daha zengin tahminler yapmak. Üçüncüsü uydu sinyali: Sentinel-2 NDVI "
      "bitki sağlığı verisini sisteme katarak arz şoklarını daha hale "
      "ulaşmadan yakalamak. Dördüncüsü sürekli streaming: NiFi ve Kafka "
      "hattını batch'in ötesine taşıyıp tam sürekli bir pipeline'a "
      "yükseltmek. Beşincisi rota analizi: Ankara-Konya gibi taşıma "
      "koridorlarını inceleyerek lojistik maliyetini fiyatlama gücünden "
      "ayırmak — hocamızın da değindiği bir nokta. Altıncısı kapsam: gerçek "
      "hal verisini 10 ilin ötesine genişletmek ve ürün kataloğunu "
      "derinleştirmek. Yani sağlam bir temel kurduk ve büyütülecek yön "
      "net.")
footer(s)

# ===========================================================================
# SLAYT 27 — REFERENCES & CLOSING
# ===========================================================================
s = new_slide(NAVY)
tf = textbox(s, 0.8, 0.55, 11.7, 0.8)
para(tf, "References & Closing", 30, WHITE, bold=True, first=True)
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3),
                         Inches(3.0), Inches(0.06))
_solid(acc, TEAL)
acc.shadow.inherit = False
refs = [
    "Peltzman, S. (2000). Prices Rise Faster Than They Fall. "
    "Journal of Political Economy, 108(3).",
    "Bacon, R. W. (1991). Rockets and feathers: asymmetric speed of "
    "adjustment of UK retail gasoline prices. Energy Economics, 13(3).",
    "Meyer, J. & von Cramon-Taubadel, S. (2004). Asymmetric Price "
    "Transmission: A Survey. Journal of Agricultural Economics, 55(3).",
    "Taylor, S. J. & Letham, B. (2018). Forecasting at Scale (Prophet). "
    "The American Statistician, 72(1).",
    "Data: marketfiyati.org.tr · İBB Tarım & Harman Hal · GDELT Project · "
    "EPİAŞ Transparency · TCMB EVDS · Open-Meteo · EPDK · Yahoo Finance.",
]
tf = textbox(s, 0.85, 1.55, 11.6, 3.0)
first = True
for r in refs:
    para(tf, r, 14, RGBColor(0xD7, 0xDE, 0xE8), bullet="▪", first=first,
         space_after=10)
    first = False
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(4.95),
                          Inches(11.6), Inches(0.04))
_solid(line, RGBColor(0x3A, 0x49, 0x63))
line.shadow.inherit = False
tf = textbox(s, 0.85, 5.2, 11.6, 1.0)
para(tf, "Thank you", 34, WHITE, bold=True, first=True)
para(tf, "Azmi Yağlı · Abdullah Zengin · Hidayet Ersin Dursun — "
         "questions welcome.", 17, TEAL, bold=True)
tf = textbox(s, 0.85, 6.55, 11.6, 0.6)
para(tf, "Disclaimer: AI assistance was used in the development and "
         "preparation of this project (AI assist kullanıldı).", 12,
     RGBColor(0x9A, 0xA6, 0xB8), italic=True, first=True)
notes(s, "27 — Kaynakça ve Kapanış",
      "Son slaytımız kaynakça ve kapanış. Asimetri analizimizin akademik "
      "temelini Peltzman, Bacon ve Meyer ile von Cramon-Taubadel'in "
      "çalışmaları oluşturuyor; tahminleme için Prophet'in orijinal "
      "makalesini referans aldık. Veri kaynaklarımızı da burada listeledik: "
      "marketfiyatı, İBB ve Harman hal, GDELT, EPİAŞ, TCMB EVDS, "
      "Open-Meteo, EPDK ve Yahoo Finance. Ayrıca şeffaflık adına bir not: "
      "projenin geliştirilmesi ve sunum hazırlığında yapay zeka "
      "araçlarından yararlandık. Toparlarsak: Türkiye gıda tedarik "
      "zincirindeki marjı uçtan uca ölçen, asimetrik fiyat geçişini "
      "sayısallaştıran, hava şoklarının yansıma hızını bulan ve tahmin "
      "üreten bütün bir büyük veri sistemi kurduk; AWS üzerinde, medallion "
      "mimarisiyle ve gerçek sonuçlarla. Bizi dinlediğiniz için teşekkür "
      "ederiz — sorularınızı memnuniyetle alırız.")
footer(s)

# ---- kaydet ----------------------------------------------------------------
out_pptx = os.path.join(HERE, "GidaRadar_Final_Presentation.pptx")
prs.save(out_pptx)
print(f"PPTX  -> {out_pptx}  ({len(prs.slides)} slayt)")

# ---- konusma_metni.md ------------------------------------------------------
md = ["# GıdaRadar — Final Sunum Konuşma Metni",
      "",
      "Slaytlar İngilizce, konuşma metni Türkçe. Her başlık bir slayta "
      "karşılık gelir; metinler `.pptx` dosyasında slayt notları olarak da "
      "gömülüdür.",
      ""]
for title, text in NOTES:
    md.append(f"## Slayt {title}")
    md.append("")
    md.append(text)
    md.append("")
out_md = os.path.join(HERE, "konusma_metni.md")
with open(out_md, "w", encoding="utf-8") as f:
    f.write("\n".join(md))
print(f"NOTES -> {out_md}  ({len(NOTES)} slayt metni)")
